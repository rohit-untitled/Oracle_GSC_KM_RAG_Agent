import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.services.extractors.pptx_extractor import (
    _extract_comment_xml_text,
    extract_presentation,
    presentation_record_to_ingestion_records,
    presentation_record_to_text,
)


class PptxExtractorTests(unittest.TestCase):
    def test_extracts_slide_text_tables_notes_and_metadata(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            path = Path(tmp_dir) / "rag_deck.pptx"
            prs = Presentation()

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Implementation Timeline"
            body = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
            body.text_frame.text = "Phase 1 readiness"

            table = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(6), Inches(1)).table
            table.cell(0, 0).text = "Milestone"
            table.cell(0, 1).text = "Owner"
            table.cell(1, 0).text = "Build | Validate"
            table.cell(1, 1).text = "KM Team"

            notes = slide.notes_slide.notes_text_frame
            notes.text = "Confirm migration dependencies."

            prs.save(path)

            record = extract_presentation(str(path))
            self.assertEqual(record.slide_count, 1)
            self.assertEqual(record.slides[0].slide_title, "Implementation Timeline")
            self.assertTrue(record.slides[0].has_tables)
            self.assertEqual(record.slides[0].skipped_media_count, 0)

            text = presentation_record_to_text(record)
            self.assertIn("Phase 1 readiness", text)
            self.assertIn("Build \\| Validate", text)
            self.assertIn("[speaker_notes]", text)
            self.assertNotIn("[embedded_object]", text)

            chunks = presentation_record_to_ingestion_records(record)
            self.assertEqual(len(chunks), 1)
            self.assertEqual(chunks[0].metadata["document_type"], "pptx")
            self.assertEqual(chunks[0].metadata["slide_number"], 1)
            self.assertIn("slide_width", chunks[0].metadata)

    def test_extract_comment_xml_text(self):
        xml = b"""
        <p:cmLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <p:cm>
                <p:text><a:r><a:t>Review the cutover date.</a:t></a:r></p:text>
            </p:cm>
        </p:cmLst>
        """

        comments = _extract_comment_xml_text(xml)
        self.assertEqual(comments, ["Review the cutover date."])


if __name__ == "__main__":
    unittest.main()
