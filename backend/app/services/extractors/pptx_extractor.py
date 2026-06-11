from __future__ import annotations

import logging
import os
import posixpath
import re
import shutil
import tempfile
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

logger = logging.getLogger(__name__)
_IMAGE_SUMMARY_CONFIG_WARNING_LOGGED = False

TITLE_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}
MAX_CHART_SERIES = 10
MAX_CHART_POINTS_PER_SERIES = 25
MEDIA_SHAPE_TYPES = {
    MSO_SHAPE_TYPE.MEDIA,
    MSO_SHAPE_TYPE.WEB_VIDEO,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
    MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
}


def _normalize_text(text: str) -> str:
    return re.sub(r"[ \t\r\f\v]+", " ", text.strip())


def _normalize_multiline_text(text: str) -> str:
    lines = [_normalize_text(line) for line in str(text or "").splitlines()]
    return "\n".join(line for line in lines if line)


@dataclass
class PptxSlideElement:
    element_type: str
    content: str
    metadata: dict[str, Any] | None = None


@dataclass
class PptxSlideRecord:
    slide_number: int
    slide_title: str | None
    elements: list[PptxSlideElement]
    speaker_notes: str | None = None
    comments: list[str] | None = None
    has_images: bool = False
    has_tables: bool = False
    has_charts: bool = False
    skipped_media_count: int = 0
    skipped_media_types: list[str] | None = None
    metadata: dict[str, Any] | None = None


@dataclass
class PptxPresentationRecord:
    file_path: str
    file_name: str
    file_type: str
    slide_count: int
    slides: list[PptxSlideRecord]
    properties: dict[str, Any] | None = None


@dataclass
class PptxIngestionRecord:
    record_id: str
    text: str
    metadata: dict[str, Any]


def _save_pptx_image_blob(blob: bytes, ext: str, image_dir: str, img_index: int) -> str:
    ext = ext if ext.startswith(".") else f".{ext}"
    img_name = f"img_{img_index}{ext or '.png'}"
    img_path = os.path.join(image_dir, img_name)
    with open(img_path, "wb") as f:
        f.write(blob)
    return img_path


def _emu_value(value: Any) -> int | None:
    try:
        if value is None:
            return None
        return int(value)
    except Exception:
        return None


def _xml_local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def _extract_comment_xml_text(xml_bytes: bytes) -> list[str]:
    try:
        root = ElementTree.fromstring(xml_bytes)
    except Exception:
        return []

    comments: list[str] = []
    for node in root.iter():
        local_name = _xml_local_name(str(node.tag)).lower()
        if local_name not in {"cm", "comment", "threadedcomment"}:
            continue

        parts: list[str] = []
        for child in node.iter():
            child_name = _xml_local_name(str(child.tag)).lower()
            if child_name not in {"text", "t"}:
                continue
            text = _normalize_multiline_text(child.text or "")
            if text:
                parts.append(text)

        comment = _normalize_multiline_text("\n".join(parts))
        if comment:
            comments.append(comment)

    return _dedupe_preserve_order(comments)


def _resolve_package_target(source_part: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def _extract_comments_by_slide_part(pptx_path: str | Path) -> dict[str, list[str]]:
    comments_by_slide: dict[str, list[str]] = {}

    try:
        with zipfile.ZipFile(pptx_path) as archive:
            names = set(archive.namelist())
            rel_names = [
                name
                for name in names
                if name.startswith("ppt/slides/_rels/")
                and name.endswith(".xml.rels")
                and name.rsplit("/", 1)[-1].startswith("slide")
            ]

            for rel_name in rel_names:
                source_part = "ppt/slides/" + rel_name.rsplit("/", 1)[-1].removesuffix(".rels")
                try:
                    rel_root = ElementTree.fromstring(archive.read(rel_name))
                except Exception:
                    continue

                slide_comments: list[str] = []
                for rel in rel_root:
                    rel_type = rel.attrib.get("Type", "")
                    target = rel.attrib.get("Target", "")
                    if not target or "comment" not in rel_type.lower():
                        continue

                    target_name = _resolve_package_target(source_part, target)
                    if target_name not in names:
                        continue
                    slide_comments.extend(_extract_comment_xml_text(archive.read(target_name)))

                if slide_comments:
                    comments_by_slide[source_part] = _dedupe_preserve_order(slide_comments)
    except zipfile.BadZipFile:
        logger.warning("PPTX file is not a valid zip package: %s", pptx_path)
    except Exception as exc:
        logger.debug("Failed to extract PPTX comments: %s", exc)

    return comments_by_slide


def _summarize_image(img_path: str) -> str:
    global _IMAGE_SUMMARY_CONFIG_WARNING_LOGGED

    image_summary_enabled = (os.getenv("IMAGE_SUMMARY_ENABLED", "true") or "true").strip().lower() in {
        "1",
        "true",
        "yes",
        "y",
    }
    if not image_summary_enabled:
        return ""

    missing_env = [
        env_name
        for env_name in ("OCI_REGION", "CONFIG_PROFILE", "COMPARTMENT_ID", "OCI_OPENAI_MODEL")
        if not os.getenv(env_name)
    ]
    if missing_env:
        if not _IMAGE_SUMMARY_CONFIG_WARNING_LOGGED:
            logger.warning(
                "PPTX image summarization is skipped because required env vars are missing: %s",
                ", ".join(missing_env),
            )
            _IMAGE_SUMMARY_CONFIG_WARNING_LOGGED = True
        return ""

    try:
        from app.services.extractors.image_summary_service import summarize_image_with_llm
    except Exception as exc:
        if not _IMAGE_SUMMARY_CONFIG_WARNING_LOGGED:
            logger.warning("Image summary service is unavailable for PPTX extraction: %s", exc)
            _IMAGE_SUMMARY_CONFIG_WARNING_LOGGED = True
        return ""

    try:
        return summarize_image_with_llm(img_path) or ""
    except Exception as exc:
        logger.warning("Failed to summarize PPTX image %s: %s", img_path, exc)
        return ""


def _shape_metadata(shape, extra: dict[str, Any] | None = None) -> dict[str, Any]:
    metadata: dict[str, Any] = {}
    try:
        metadata["shape_name"] = shape.name
    except Exception:
        pass

    try:
        metadata["shape_type"] = str(shape.shape_type)
    except Exception:
        pass

    for attr in ("left", "top", "width", "height"):
        try:
            value = _emu_value(getattr(shape, attr))
        except Exception:
            value = None
        if value is not None:
            metadata[attr] = value

    try:
        rotation = getattr(shape, "rotation", None)
        if rotation not in (None, 0):
            metadata["rotation"] = rotation
    except Exception:
        pass

    try:
        if shape.is_placeholder:
            metadata["is_placeholder"] = True
            metadata["placeholder_type"] = str(shape.placeholder_format.type)
            metadata["placeholder_index"] = shape.placeholder_format.idx
    except Exception:
        pass

    try:
        c_nv_pr = shape.element.xpath(".//p:cNvPr")
        if c_nv_pr:
            title = c_nv_pr[0].get("title")
            description = c_nv_pr[0].get("descr")
            if title:
                metadata["alt_title"] = title
            if description:
                metadata["alt_text"] = description
    except Exception:
        pass

    try:
        address = shape.click_action.hyperlink.address
        if address:
            metadata["shape_hyperlink"] = address
    except Exception:
        pass

    if extra:
        metadata.update({key: value for key, value in extra.items() if value not in (None, "", [])})
    return metadata or None


def _shape_hyperlink_line(metadata: dict[str, Any] | None) -> str | None:
    if not metadata:
        return None
    link = metadata.get("shape_hyperlink")
    if not link:
        return None
    return f"Link: {link}"


def _is_title_shape(shape) -> bool:
    try:
        if not shape.is_placeholder:
            return False
        return shape.placeholder_format.type in TITLE_PLACEHOLDER_TYPES
    except Exception:
        return False


def _paragraph_text_and_links(paragraph) -> tuple[str, list[str]]:
    links: list[str] = []
    run_texts: list[str] = []

    try:
        for run in paragraph.runs:
            text = run.text or ""
            if text:
                run_texts.append(text)
            try:
                address = run.hyperlink.address
                if address:
                    links.append(address)
            except Exception:
                pass
    except Exception:
        pass

    text = "".join(run_texts).strip() or (paragraph.text or "").strip()
    text = _normalize_text(text)
    if not text:
        return "", links

    for link in links:
        if link not in text:
            text = f"{text} (link: {link})"
    return text, links


def _extract_text_from_text_frame(text_frame) -> tuple[list[str], list[str]]:
    texts: list[str] = []
    links: list[str] = []

    try:
        for paragraph in text_frame.paragraphs:
            text, paragraph_links = _paragraph_text_and_links(paragraph)
            if not text:
                continue
            links.extend(paragraph_links)
            if getattr(paragraph, "level", 0) and paragraph.level > 0:
                texts.append(f"- {text}")
            else:
                texts.append(text)
    except Exception as exc:
        logger.debug("Failed to extract text frame: %s", exc)

    return texts, _dedupe_preserve_order(links)


def _extract_xml_text(shape) -> list[str]:
    try:
        texts = [
            _normalize_text(node.text)
            for node in shape.element.xpath(".//a:t")
            if getattr(node, "text", None) and _normalize_text(node.text)
        ]
    except Exception:
        return []
    return _dedupe_preserve_order(texts)


def _dedupe_preserve_order(values: list[Any]) -> list[Any]:
    seen: set[str] = set()
    deduped: list[Any] = []
    for value in values:
        key = str(value)
        if key in seen:
            continue
        seen.add(key)
        deduped.append(value)
    return deduped


def _table_cell_text(cell) -> str:
    try:
        texts, _links = _extract_text_from_text_frame(cell.text_frame)
        text = " ".join(texts) if texts else cell.text_frame.text
    except Exception:
        text = ""
    return _normalize_text(text.replace("\n", " "))


def _escape_markdown_table_cell(value: str) -> str:
    return value.replace("\\", "\\\\").replace("|", "\\|")


def _extract_table_data(shape) -> str | None:
    try:
        table = shape.table
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([_table_cell_text(cell) for cell in row.cells])
    except Exception as exc:
        logger.debug("Failed to extract PPTX table: %s", exc)
        return None

    rows = [row for row in rows if any(cell for cell in row)]
    if not rows:
        return None

    col_count = max(len(row) for row in rows)
    normalized_rows = [
        [_escape_markdown_table_cell(cell) for cell in row + [""] * (col_count - len(row))]
        for row in rows
    ]
    markdown_rows = ["| " + " | ".join(row) + " |" for row in normalized_rows]

    if len(normalized_rows) > 1 and any(cell for cell in normalized_rows[0]):
        separator = "| " + " | ".join("---" for _ in range(col_count)) + " |"
        markdown_rows.insert(1, separator)

    return "\n".join(markdown_rows)


def _chart_values_from_cache(parent, tag_name: str) -> list[str]:
    values: list[str] = []
    try:
        containers = parent.xpath(f".//c:{tag_name}")
    except Exception:
        return values

    for container in containers:
        try:
            values.extend(
                _normalize_text(node.text)
                for node in container.xpath(".//c:pt/c:v")
                if getattr(node, "text", None) and _normalize_text(node.text)
            )
        except Exception:
            continue
    return values


def _chart_series_name(series_node, fallback: str) -> str:
    try:
        node_values = [
            _normalize_text(node.text)
            for node in series_node.xpath(".//c:tx//c:v")
            if getattr(node, "text", None) and _normalize_text(node.text)
        ]
        if node_values:
            return node_values[0]
    except Exception:
        pass
    return fallback


def _format_chart_points(labels: list[str], values: list[str], max_points: int) -> str:
    if not values:
        return ""

    formatted: list[str] = []
    for idx, value in enumerate(values[:max_points]):
        label = labels[idx] if idx < len(labels) and labels[idx] else f"point_{idx + 1}"
        formatted.append(f"{label}={value}")

    remaining = len(values) - max_points
    suffix = f"; ... {remaining} more" if remaining > 0 else ""
    return "; ".join(formatted) + suffix


def _extract_chart_data_from_xml(shape) -> list[str]:
    try:
        chart = shape.chart
        series_nodes = chart.element.xpath(".//c:ser")
    except Exception:
        return []

    parts: list[str] = []
    for idx, series_node in enumerate(series_nodes[:MAX_CHART_SERIES], start=1):
        name = _chart_series_name(series_node, f"Series {idx}")
        categories = _chart_values_from_cache(series_node, "cat")
        values = _chart_values_from_cache(series_node, "val")
        x_values = _chart_values_from_cache(series_node, "xVal")
        y_values = _chart_values_from_cache(series_node, "yVal")
        bubble_sizes = _chart_values_from_cache(series_node, "bubbleSize")

        if values:
            point_text = _format_chart_points(categories, values, MAX_CHART_POINTS_PER_SERIES)
            parts.append(f"{name}: {point_text}")
            continue

        if x_values or y_values:
            points: list[str] = []
            point_count = max(len(x_values), len(y_values), len(bubble_sizes))
            for point_idx in range(min(point_count, MAX_CHART_POINTS_PER_SERIES)):
                x_value = x_values[point_idx] if point_idx < len(x_values) else ""
                y_value = y_values[point_idx] if point_idx < len(y_values) else ""
                bubble_size = bubble_sizes[point_idx] if point_idx < len(bubble_sizes) else ""
                point = f"x={x_value}, y={y_value}"
                if bubble_size:
                    point = f"{point}, size={bubble_size}"
                points.append(point)
            remaining = point_count - MAX_CHART_POINTS_PER_SERIES
            suffix = f"; ... {remaining} more" if remaining > 0 else ""
            parts.append(f"{name}: " + "; ".join(points) + suffix)

    if len(series_nodes) > MAX_CHART_SERIES:
        parts.append(f"... {len(series_nodes) - MAX_CHART_SERIES} more series")

    return parts


def _extract_chart_data_from_api(shape) -> list[str]:
    try:
        chart = shape.chart
    except Exception:
        return []

    parts: list[str] = []
    try:
        for plot in chart.plots:
            categories = []
            try:
                categories = [str(category) for category in plot.categories]
            except Exception:
                pass

            for series in list(plot.series)[:MAX_CHART_SERIES]:
                try:
                    values = [str(value) for value in series.values]
                except Exception:
                    values = []
                if not values:
                    continue
                name = getattr(series, "name", None) or f"Series {getattr(series, 'index', 0) + 1}"
                point_text = _format_chart_points(categories, values, MAX_CHART_POINTS_PER_SERIES)
                parts.append(f"{name}: {point_text}")
    except Exception:
        return []

    return parts


def _extract_chart_data(shape) -> str | None:
    try:
        if not getattr(shape, "has_chart", False):
            return None

        chart = shape.chart
        parts = [f"Chart type: {chart.chart_type}"]

        try:
            if chart.has_title and chart.chart_title.text_frame:
                title = _normalize_multiline_text(chart.chart_title.text_frame.text)
                if title:
                    parts.append(f"Title: {title}")
        except Exception:
            pass

        data_parts = _extract_chart_data_from_xml(shape) or _extract_chart_data_from_api(shape)
        parts.extend(data_parts)

        return "\n".join(_dedupe_preserve_order(parts)) if len(parts) > 1 else parts[0]
    except Exception as exc:
        logger.warning("Failed to extract PPTX chart: %s", exc)
        return None


def _extract_speaker_notes(slide) -> str | None:
    try:
        if not slide.has_notes_slide:
            return None
        notes_text_frame = slide.notes_slide.notes_text_frame
        if not notes_text_frame:
            return None
        notes_text = _normalize_multiline_text(notes_text_frame.text)
        return notes_text or None
    except Exception as exc:
        logger.debug("Failed to extract speaker notes: %s", exc)
        return None


def _extract_core_properties(prs) -> dict[str, Any] | None:
    try:
        props = prs.core_properties
    except Exception:
        return None

    metadata: dict[str, Any] = {}
    for attr in (
        "author",
        "category",
        "comments",
        "content_status",
        "created",
        "identifier",
        "keywords",
        "language",
        "last_modified_by",
        "last_printed",
        "modified",
        "revision",
        "subject",
        "title",
        "version",
    ):
        try:
            value = getattr(props, attr)
        except Exception:
            continue
        if value is None or value == "":
            continue
        metadata[attr] = value.isoformat() if hasattr(value, "isoformat") else value
    return metadata or None


def _slide_part_name(slide) -> str | None:
    try:
        return str(slide.part.partname).lstrip("/")
    except Exception:
        return None


def _slide_metadata(prs, slide) -> dict[str, Any]:
    metadata: dict[str, Any] = {}
    try:
        metadata["slide_width"] = int(prs.slide_width)
        metadata["slide_height"] = int(prs.slide_height)
    except Exception:
        pass

    try:
        show = slide.element.get("show")
        if show is not None:
            metadata["is_hidden"] = show == "0"
    except Exception:
        pass

    try:
        part_name = _slide_part_name(slide)
        if part_name:
            metadata["slide_part_name"] = part_name
    except Exception:
        pass

    try:
        if slide.slide_layout and slide.slide_layout.name:
            metadata["slide_layout"] = slide.slide_layout.name
    except Exception:
        pass

    return metadata


def extract_presentation(pptx_path: str, image_dir: str | None = None) -> PptxPresentationRecord:
    prs = Presentation(pptx_path)
    path = Path(pptx_path)
    temp_dir_created = image_dir is None
    temp_dir = image_dir or tempfile.mkdtemp(prefix="pptx_extract_")
    os.makedirs(temp_dir, exist_ok=True)
    comments_by_slide = _extract_comments_by_slide_part(path)

    slides: list[PptxSlideRecord] = []
    img_counter = 0

    try:
        for slide_no, slide in enumerate(prs.slides, start=1):
            elements: list[PptxSlideElement] = []
            slide_title = None
            has_images = False
            has_tables = False
            has_charts = False
            skipped_media_count = 0
            skipped_media_types: list[str] = []
            slide_metadata = _slide_metadata(prs, slide)
            slide_part_name = slide_metadata.get("slide_part_name")
            slide_comments = comments_by_slide.get(slide_part_name, []) if slide_part_name else []

            def process_shape(shape, group_path: list[str] | None = None, z_order: int | None = None) -> None:
                nonlocal img_counter, slide_title, has_images, has_tables, has_charts, skipped_media_count

                group_path = group_path or []
                metadata_extra: dict[str, Any] = {}
                if group_path:
                    metadata_extra["group_path"] = " > ".join(group_path)
                if z_order is not None:
                    metadata_extra["z_order"] = z_order

                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                    try:
                        group_name = getattr(shape, "name", "Group")
                    except Exception:
                        group_name = "Group"
                    for child_idx, child in enumerate(shape.shapes, start=1):
                        try:
                            process_shape(child, [*group_path, group_name], child_idx)
                        except Exception as exc:
                            logger.warning("Failed to extract grouped shape on slide %s: %s", slide_no, exc)
                    return

                if getattr(shape, "shape_type", None) in MEDIA_SHAPE_TYPES:
                    skipped_media_count += 1
                    media_type = str(getattr(shape, "shape_type", "unknown"))
                    if media_type not in skipped_media_types:
                        skipped_media_types.append(media_type)
                    logger.debug("Skipped PPTX media object on slide %s: %s", slide_no, media_type)
                    return

                if getattr(shape, "shape_type", None) in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}:
                    has_images = True
                    try:
                        image = shape.image
                        img_path = _save_pptx_image_blob(image.blob, image.ext or "png", temp_dir, img_counter)
                        image_index = img_counter
                        img_counter += 1
                    except Exception as exc:
                        logger.warning("Failed to save PPTX image on slide %s: %s", slide_no, exc)
                        image_index = img_counter
                        img_path = None

                    metadata = _shape_metadata(shape, {**metadata_extra, "image_index": image_index})
                    content_parts: list[str] = []
                    if metadata:
                        alt_title = metadata.get("alt_title")
                        alt_text = metadata.get("alt_text")
                        if alt_title:
                            content_parts.append(f"Image title: {alt_title}")
                        if alt_text:
                            content_parts.append(f"Image alt text: {alt_text}")
                        hyperlink_line = _shape_hyperlink_line(metadata)
                        if hyperlink_line:
                            content_parts.append(hyperlink_line)
                    if img_path:
                        summary_text = _normalize_multiline_text(_summarize_image(img_path))
                        if summary_text and summary_text.upper() != "EMPTY":
                            content_parts.append(f"Image summary: {summary_text}")
                    if content_parts:
                        elements.append(
                            PptxSlideElement(
                                element_type="image",
                                content="\n".join(content_parts),
                                metadata=metadata,
                            )
                        )
                    return

                if getattr(shape, "has_table", False):
                    has_tables = True
                    table_content = _extract_table_data(shape)
                    if table_content:
                        metadata = _shape_metadata(shape, metadata_extra)
                        hyperlink_line = _shape_hyperlink_line(metadata)
                        if hyperlink_line:
                            table_content = f"{table_content}\n{hyperlink_line}"
                        elements.append(
                            PptxSlideElement(
                                element_type="table",
                                content=table_content,
                                metadata=metadata,
                            )
                        )
                    return

                if getattr(shape, "has_chart", False) or getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.CHART:
                    has_charts = True
                    chart_text = _extract_chart_data(shape)
                    if chart_text:
                        metadata = _shape_metadata(shape, metadata_extra)
                        hyperlink_line = _shape_hyperlink_line(metadata)
                        if hyperlink_line:
                            chart_text = f"{chart_text}\n{hyperlink_line}"
                        elements.append(
                            PptxSlideElement(
                                element_type="chart",
                                content=chart_text,
                                metadata=metadata,
                            )
                        )
                    return

                if getattr(shape, "has_text_frame", False):
                    texts, links = _extract_text_from_text_frame(shape.text_frame)
                    if texts:
                        text_content = "\n".join(texts)
                        metadata = _shape_metadata(shape, {**metadata_extra, "hyperlinks": links})
                        hyperlink_line = _shape_hyperlink_line(metadata)
                        if hyperlink_line and hyperlink_line not in text_content:
                            text_content = f"{text_content}\n{hyperlink_line}"
                        is_title = _is_title_shape(shape)
                        if is_title and not slide_title:
                            slide_title = _normalize_multiline_text(text_content)
                        elements.append(
                            PptxSlideElement(
                                element_type="title" if is_title else "text",
                                content=text_content,
                                metadata=metadata,
                            )
                        )
                        return

                xml_texts = _extract_xml_text(shape)
                if xml_texts:
                    metadata = _shape_metadata(shape, metadata_extra)
                    hyperlink_line = _shape_hyperlink_line(metadata)
                    if hyperlink_line:
                        xml_texts.append(hyperlink_line)
                    elements.append(
                        PptxSlideElement(
                            element_type="embedded_text",
                            content="\n".join(xml_texts),
                            metadata=metadata,
                        )
                    )

            for shape_idx, shape in enumerate(slide.shapes, start=1):
                try:
                    process_shape(shape, z_order=shape_idx)
                except Exception as exc:
                    logger.warning("Failed to extract shape on slide %s: %s", slide_no, exc)

            speaker_notes = _extract_speaker_notes(slide)
            if not slide_title:
                for element in elements:
                    if element.element_type in {"title", "text"} and element.content:
                        slide_title = _normalize_multiline_text(element.content.splitlines()[0])
                        break

            slides.append(
                PptxSlideRecord(
                    slide_number=slide_no,
                    slide_title=slide_title,
                    elements=elements,
                    speaker_notes=speaker_notes,
                    comments=slide_comments or None,
                    has_images=has_images,
                    has_tables=has_tables,
                    has_charts=has_charts,
                    skipped_media_count=skipped_media_count,
                    skipped_media_types=skipped_media_types or None,
                    metadata=slide_metadata or None,
                )
            )

        return PptxPresentationRecord(
            file_path=str(path.resolve()),
            file_name=path.name,
            file_type=path.suffix.lstrip("."),
            slide_count=len(slides),
            slides=slides,
            properties=_extract_core_properties(prs),
        )

    finally:
        if temp_dir_created:
            shutil.rmtree(temp_dir, ignore_errors=True)


def presentation_record_to_dict(record: PptxPresentationRecord) -> dict[str, Any]:
    return asdict(record)


def _slide_element_to_text(element: PptxSlideElement) -> str:
    content = _normalize_multiline_text(element.content)
    if not content:
        return ""
    return f"[{element.element_type}]\n{content}"


def presentation_record_to_ingestion_records(record: PptxPresentationRecord) -> list[PptxIngestionRecord]:
    ingestion_records: list[PptxIngestionRecord] = []

    for slide in record.slides:
        element_texts: list[str] = []
        element_types: list[str] = []

        for element in slide.elements:
            if element.element_type not in element_types:
                element_types.append(element.element_type)
            element_text = _slide_element_to_text(element)
            if element_text:
                element_texts.append(element_text)

        slide_content_parts = [
            f"Presentation: {record.file_name}",
            f"Slide {slide.slide_number}" + (f": {slide.slide_title}" if slide.slide_title else ""),
        ]
        slide_content_parts.extend(element_texts)
        if slide.speaker_notes:
            slide_content_parts.append(f"[speaker_notes]\n{_normalize_multiline_text(slide.speaker_notes)}")
        if slide.comments:
            comments_text = "\n".join(f"- {_normalize_multiline_text(comment)}" for comment in slide.comments)
            slide_content_parts.append(f"[comments]\n{comments_text}")

        if not element_texts and not slide.speaker_notes and not slide.comments and not slide.slide_title:
            continue

        slide_content = "\n\n".join(part for part in slide_content_parts if part)
        if not slide_content.strip():
            continue

        metadata = {
            "document_type": "pptx",
            "file_name": record.file_name,
            "file_path": record.file_path,
            "file_type": record.file_type,
            "slide_number": slide.slide_number,
            "slide_title": slide.slide_title,
            "element_types": element_types,
            "has_speaker_notes": slide.speaker_notes is not None,
            "has_images": slide.has_images,
            "has_tables": slide.has_tables,
            "has_charts": slide.has_charts,
            "has_comments": bool(slide.comments),
            "skipped_media_count": slide.skipped_media_count,
            "skipped_media_types": slide.skipped_media_types or [],
            "total_elements": len(slide.elements),
            "chunk_type": "pptx_slide",
        }
        if slide.metadata:
            metadata.update(slide.metadata)
        if record.properties:
            metadata["presentation_properties"] = record.properties

        ingestion_records.append(
            PptxIngestionRecord(
                record_id=f"{record.file_name}:slide_{slide.slide_number}",
                text=slide_content,
                metadata=metadata,
            )
        )

    return ingestion_records


def presentation_record_to_text(record: PptxPresentationRecord) -> str:
    lines = [
        f"Presentation: {record.file_name}",
        f"Path: {record.file_path}",
        f"Type: {record.file_type}",
        f"Slides: {record.slide_count}",
    ]

    if record.properties:
        property_text = ", ".join(f"{key}: {value}" for key, value in record.properties.items())
        lines.append(f"Properties: {property_text}")

    for slide in record.slides:
        lines.append("")
        lines.append(f"## Slide {slide.slide_number}")
        if slide.slide_title:
            lines.append(f"Title: {slide.slide_title}")

        for element in slide.elements:
            element_text = _slide_element_to_text(element)
            if element_text:
                lines.append("")
                lines.append(element_text)

        if slide.speaker_notes:
            lines.append("")
            lines.append(f"[speaker_notes]\n{_normalize_multiline_text(slide.speaker_notes)}")
        if slide.comments:
            lines.append("")
            lines.append("[comments]")
            lines.extend(f"- {_normalize_multiline_text(comment)}" for comment in slide.comments)

    return "\n".join(lines)


def ingestion_records_to_text(records: list[PptxIngestionRecord]) -> str:
    lines = [f"Ingestion records: {len(records)}"]

    for record in records:
        metadata = record.metadata
        lines.append("")
        lines.append(f"## {record.record_id}")
        lines.append(f"Slide {metadata.get('slide_number')}: {metadata.get('slide_title', 'Untitled')}")
        lines.append(f"Elements: {', '.join(metadata.get('element_types', []))}")
        lines.append(f"Content: {record.text[:500]}...")

    return "\n".join(lines)


def ingestion_records_to_dict(records: list[PptxIngestionRecord]) -> list[dict[str, Any]]:
    return [asdict(record) for record in records]


def extract_pptx_with_formatting_in_sequence(pptx_path: str, image_dir: str | None = None) -> str:
    record = extract_presentation(pptx_path, image_dir)
    return presentation_record_to_text(record)
